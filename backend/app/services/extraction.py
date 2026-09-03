"""文本抽取服务。

输入文件字节 + 文件名 / mime，输出抽取全文草稿 + 状态。纯 Python 抽取库
（txt/md 直读、pdf 用 pypdf、docx 用 python-docx、pptx 用 python-pptx、
xlsx 用 openpyxl 转 markdown 表格）。本模块只识别需 OCR 的页，OCR 由独立本地引擎执行。

边界：不接真实 LLM / 大模型；不支持 旧版 .xls
（标 `unsupported`，不崩溃、不阻断任务创建）；不做切块 / 向量化。

安全：抽取全文是**用户业务内容**，可能含 `s3://` / `internal://` / URL 等字样——
它只准进草稿 / 资产侧，**绝不准进审计 extra/after**（审计只放安全元数据）。
"""

from __future__ import annotations

import io
import logging
import pickle
import subprocess
import sys
import zipfile
from dataclasses import dataclass

from app.core.logging import safe_log_exception
from app.core.text_safety import EXTRACTED_TEXT_MAX_CHARS, SafetyStats, sanitize_text

_logger = logging.getLogger(__name__)

# 抽取草稿全文上限（防超大文本放大）；超过则截断。
MAX_EXTRACT_CHARS = EXTRACTED_TEXT_MAX_CHARS

# All parsers below consume attacker-controlled bytes.  The limits are deliberately
# independent from the HTTP upload limit: a small ZIP can expand into a very large
# Office document, and a valid-looking PDF can still make a parser loop forever.
MAX_ARCHIVE_ENTRIES = 20_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
MAX_ARCHIVE_COMPRESSION_RATIO = 100
EXTRACTION_TIMEOUT_SECONDS = 45.0
MAX_PDF_PAGES = 2_000
MAX_PPTX_SLIDES = 1_000
MAX_XLSX_SHEETS = 200
MAX_XLSX_ROWS = 200_000
MAX_XLSX_COLUMNS = 2_000

# 直接按文本读取的扩展名。
_TEXT_EXT = {"txt", "md", "markdown", "csv", "log", "text", "rst"}

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp"}
_KNOWN_EXTENSIONS = (
    _TEXT_EXT
    | _IMAGE_EXTENSIONS
    | {
        "pdf",
        "doc",
        "docx",
        "ppt",
        "pptx",
        "xls",
        "xlsx",
    }
)


@dataclass(frozen=True)
class ExtractionPage:
    page_number: int
    text: str
    status: str  # extracted | ocr_required


@dataclass(frozen=True)
class ExtractionResult:
    """抽取结果。status ∈ extracted / ocr_required / unsupported / failed / empty。"""

    text: str
    status: str
    error_type: str | None
    error_message: str | None
    char_count: int
    pages: tuple[ExtractionPage, ...] = ()
    source_kind: str = "document"
    safety_stats: SafetyStats = SafetyStats()


def _ext(file_name: str | None) -> str:
    name = file_name or ""
    return name.rsplit(".", 1)[1].lower() if "." in name else ""


def _extract_pdf(content: bytes) -> tuple[str, tuple[ExtractionPage, ...]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    if len(reader.pages) > MAX_PDF_PAGES:
        raise _ControlledExtractionError(
            "extraction_structure_limit", "PDF 页数超过安全处理上限，请拆分后重新上传。"
        )
    parts: list[str] = []
    pages: list[ExtractionPage] = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        # 页码标记（D1 阶段3）：供 chunk 注册表记录 source_page / 查看原文定位。
        clean = page_text.strip()
        pages.append(ExtractionPage(index, clean, "extracted" if clean else "ocr_required"))
        parts.append(f"{{{{page:{index}}}}}\n{clean}" if clean else "")
    return "\n".join(parts), tuple(pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def _pptx_shape_text(shape) -> list[str]:
    """Extract readable text from one shape while preserving its child order."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        parts: list[str] = []
        for child in shape.shapes:
            parts.extend(_pptx_shape_text(child))
        return parts
    if getattr(shape, "has_table", False):
        rows: list[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        return rows
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        return [text] if text else []
    return []


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    if len(presentation.slides) > MAX_PPTX_SLIDES:
        raise _ControlledExtractionError(
            "extraction_structure_limit", "幻灯片数量超过安全处理上限，请拆分后重新上传。"
        )
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_pptx_shape_text(shape))
        if parts:
            slides.append(f"[幻灯片 {index}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def _xlsx_cell(value: object) -> str:
    """单元格 → 单行文本（`|` / 换行转义，避免破坏 markdown 表格）。"""
    if value is None:
        return ""
    return str(value).strip().replace("|", "\\|").replace("\n", " ").replace("\r", " ")


def _extract_xlsx(content: bytes) -> str:
    """xlsx → 每个 sheet 一个 markdown 表格（sheet 名 + 表头 + 行）。

    read_only + data_only：大文件低内存；公式取 Excel 缓存值（无缓存值则为空）。
    首行视为表头（空表头回退列号 A/B/…），其余为数据行。
    """
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheets: list[str] = []
    try:
        if len(workbook.worksheets) > MAX_XLSX_SHEETS:
            raise _ControlledExtractionError(
                "extraction_structure_limit", "工作表数量超过安全处理上限，请拆分后重新上传。"
            )
        total_rows = 0
        for sheet in workbook.worksheets:
            rows: list[list[object]] = []
            for row in sheet.iter_rows(values_only=True):
                total_rows += 1
                if total_rows > MAX_XLSX_ROWS or len(row) > MAX_XLSX_COLUMNS:
                    raise _ControlledExtractionError(
                        "extraction_structure_limit",
                        "表格行列规模超过安全处理上限，请拆分或精简后重新上传。",
                    )
                if any(cell is not None and str(cell).strip() for cell in row):
                    rows.append(list(row))
            if not rows:
                continue
            width = max(len(row) for row in rows)
            first = rows[0]
            header = [_xlsx_cell(first[i] if i < len(first) else "") for i in range(width)]
            header = [h or (chr(65 + i) if i < 26 else f"列{i + 1}") for i, h in enumerate(header)]
            lines = [
                f"## {sheet.title or 'Sheet'}",
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(["---"] * width) + " |",
            ]
            for row in rows[1:]:
                cells = [_xlsx_cell(row[i] if i < len(row) else "") for i in range(width)]
                lines.append("| " + " | ".join(cells) + " |")
            sheets.append("\n".join(lines))
    finally:
        workbook.close()
    return "\n\n".join(sheets)


class _ControlledExtractionError(Exception):
    """A safe, user-actionable rejection raised before a parser sees the bytes."""

    def __init__(self, code: str, message: str) -> None:
        self.code = code
        self.message = message
        super().__init__(code)


def _office_package_kind(content: bytes, ext: str) -> None:
    """Validate OOXML ZIP structure and bound decompression before library parsing."""
    if content.startswith(b"\xd0\xcf\x11\xe0"):
        raise _ControlledExtractionError(
            "extraction_password_protected",
            "该 Office 文件已加密或需要密码保护；平台不会接收或请求密码，请移除保护后重新上传。",
        )
    if not content.startswith(b"PK\x03\x04"):
        raise _ControlledExtractionError(
            "extraction_format_mismatch", "文件内容与扩展名不匹配，请确认文件类型后重新上传。"
        )
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as package:
            infos = package.infolist()
            if len(infos) > MAX_ARCHIVE_ENTRIES:
                raise _ControlledExtractionError(
                    "extraction_structure_limit",
                    "文件内部结构过于复杂，无法安全处理，请精简后重新上传。",
                )
            total = sum(info.file_size for info in infos)
            compressed = sum(info.compress_size for info in infos)
            if total > MAX_ARCHIVE_UNCOMPRESSED_BYTES or (
                compressed and total / compressed > MAX_ARCHIVE_COMPRESSION_RATIO
            ):
                raise _ControlledExtractionError(
                    "extraction_archive_limit",
                    "文件解压后体积异常，无法安全处理，请使用原始或精简后的文件。",
                )
            names = set(package.namelist())
            if {"EncryptionInfo", "EncryptedPackage"}.intersection(names):
                raise _ControlledExtractionError(
                    "extraction_password_protected",
                    "该 Office 文件已加密或需要密码保护；平台不会接收或请求密码，请移除保护后重新上传。",
                )
            required = {
                "docx": "word/document.xml",
                "pptx": "ppt/presentation.xml",
                "xlsx": "xl/workbook.xml",
            }[ext]
            if "[Content_Types].xml" not in names or required not in names:
                raise _ControlledExtractionError(
                    "extraction_format_mismatch",
                    "文件内容与扩展名不匹配，请确认文件类型后重新上传。",
                )
    except _ControlledExtractionError:
        raise
    except (zipfile.BadZipFile, OSError, ValueError) as exc:
        raise _ControlledExtractionError(
            "extraction_corrupt", "文件内容已损坏，无法安全解析，请重新导出或上传原始文件。"
        ) from exc


def _validate_pdf(content: bytes) -> None:
    if not content.startswith(b"%PDF-"):
        raise _ControlledExtractionError(
            "extraction_format_mismatch", "文件内容与扩展名不匹配，请确认文件类型后重新上传。"
        )
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        if reader.is_encrypted:
            raise _ControlledExtractionError(
                "extraction_password_protected",
                "该 PDF 已加密或需要密码保护；平台不会接收或请求密码，请移除保护后重新上传。",
            )
    except _ControlledExtractionError:
        raise
    except Exception as exc:  # parser-specific errors are intentionally not surfaced
        raise _ControlledExtractionError(
            "extraction_corrupt", "文件内容已损坏，无法安全解析，请重新导出或上传原始文件。"
        ) from exc


def _validate_image(content: bytes) -> None:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(content)) as image:
            image.verify()
    except Exception as exc:
        raise _ControlledExtractionError(
            "extraction_corrupt", "图片内容已损坏或格式不正确，请重新导出后上传。"
        ) from exc


def _extract_unbounded(
    content: bytes, *, file_name: str | None, mime: str | None
) -> ExtractionResult:
    """Parser routing executed only in the isolated child process."""
    ext = _ext(file_name)
    mime = (mime or "").lower()
    pages: tuple[ExtractionPage, ...] = ()
    mime_fallback = not ext or ext not in _KNOWN_EXTENSIONS
    if not content:
        return ExtractionResult(
            "", "empty", "extraction_empty", "文件为空，请选择包含内容的文件后重试。", 0
        )
    if ext in _TEXT_EXT or (mime_fallback and mime.startswith("text/")):
        if content.startswith((b"%PDF-", b"PK\x03\x04", b"\xd0\xcf\x11\xe0")):
            raise _ControlledExtractionError(
                "extraction_format_mismatch", "文件内容与扩展名不匹配，请确认文件类型后重新上传。"
            )
        text = content.decode("utf-8", errors="replace")
    elif ext == "pdf" or (mime_fallback and mime == "application/pdf"):
        _validate_pdf(content)
        text, pages = _extract_pdf(content)
    elif ext == "docx" or (mime_fallback and mime == _DOCX_MIME):
        _office_package_kind(content, "docx")
        text = _extract_docx(content)
    elif ext == "pptx" or (mime_fallback and mime == _PPTX_MIME):
        _office_package_kind(content, "pptx")
        text = _extract_pptx(content)
    elif ext == "xlsx" or (mime_fallback and mime == _XLSX_MIME):
        _office_package_kind(content, "xlsx")
        text = _extract_xlsx(content)
    elif ext in _IMAGE_EXTENSIONS or (mime_fallback and mime.startswith("image/")):
        _validate_image(content)
        return ExtractionResult(
            "", "ocr_required", None, None, 0, (ExtractionPage(1, "", "ocr_required"),), "image"
        )
    else:
        unsupported_message = (
            "当前 .ppt 格式暂不支持自动提取（文件已落盘，请人工补全内容）"
            if ext == "ppt"
            else "旧版 .xls 暂不支持自动提取（文件已落盘），请另存为 .xlsx 后重新上传"
            if ext == "xls"
            else f"暂不支持从 .{ext or '该类型'} 文件抽取文本（文件已落盘，请人工补全内容）"
        )
        return ExtractionResult("", "unsupported", "extraction_unsupported", unsupported_message, 0)
    safe_text = sanitize_text(text, max_chars=MAX_EXTRACT_CHARS)
    text = safe_text.value.strip()
    if pages:
        pages = tuple(
            ExtractionPage(
                page.page_number,
                sanitize_text(page.text, max_chars=MAX_EXTRACT_CHARS).value,
                page.status,
            )
            for page in pages
        )
    is_pdf = ext == "pdf" or (mime_fallback and mime == "application/pdf")
    if is_pdf and any(page.status == "ocr_required" for page in pages):
        return ExtractionResult(
            text, "ocr_required", None, None, len(text), pages, "pdf", safe_text.stats
        )
    if not text:
        return ExtractionResult(
            "",
            "empty",
            "extraction_empty",
            "未能从文件内容中抽取到文本（可能是扫描件 / 纯图片），请人工补全",
            0,
        )
    return ExtractionResult(
        text,
        "extracted",
        None,
        None,
        len(text),
        pages if is_pdf else (),
        "pdf" if is_pdf else "document",
        safe_text.stats,
    )


def extract_text(content: bytes, *, file_name: str | None, mime: str | None) -> ExtractionResult:
    """Parse user bytes in a killable child process with bounded package complexity."""
    # subprocess (rather than multiprocessing) is intentional: a Celery prefork
    # worker can itself be daemonised and is then forbidden from creating another
    # multiprocessing child, while it may safely launch and kill a subprocess.
    try:
        child = subprocess.Popen(
            [sys.executable, "-m", "app.services.extraction_worker"],
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
        )
        try:
            output, _ = child.communicate(
                pickle.dumps((content, file_name, mime), protocol=pickle.HIGHEST_PROTOCOL),
                timeout=EXTRACTION_TIMEOUT_SECONDS,
            )
        except subprocess.TimeoutExpired:
            child.kill()
            child.communicate()
            return ExtractionResult(
                "",
                "failed",
                "extraction_timeout",
                "文件解析超过安全时限，未处理任何密码；请精简或重新导出后重试。",
                0,
            )
        try:
            kind, payload = pickle.loads(output)
        except (EOFError, pickle.UnpicklingError, ValueError, TypeError):
            return ExtractionResult(
                "",
                "failed",
                "extraction_failed",
                "文件内容无法解析（可能已损坏或与扩展名不符），请重新上传或人工补全",
                0,
            )
        if kind == "result" and isinstance(payload, ExtractionResult):
            return payload
        if kind == "controlled":
            code, message = payload
            return ExtractionResult("", "failed", code, message, 0)
        return ExtractionResult(
            "",
            "failed",
            "extraction_failed",
            "文件内容无法解析（可能已损坏或与扩展名不符），请重新上传或人工补全",
            0,
        )
    except Exception as exc:  # subprocess startup failure is also safely terminalized
        safe_log_exception(
            _logger, "extraction_failed", exc, include_summary=False, level=logging.WARNING
        )
        return ExtractionResult(
            text="",
            status="failed",
            error_type="extraction_failed",
            error_message="文件内容无法解析（可能已损坏或与扩展名不符），请重新上传或人工补全",
            char_count=0,
        )
