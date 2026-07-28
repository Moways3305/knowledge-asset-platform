"""文本抽取服务单测（IMPLEMENT-14）。

覆盖 txt/md 直读、PDF（pypdf）、docx（python-docx）、pptx（python-pptx）真实抽取，
以及 unsupported / empty / failed 分支。不依赖数据库。
"""

from __future__ import annotations

import io
import logging

from app.services.extraction import MAX_EXTRACT_CHARS, extract_text


def _make_pdf(text: str) -> bytes:
    """构造一个最小但合法的 PDF（含一段可抽取文本），用于真实 PDF 抽取测试。"""
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
    ]
    stream = b"BT /F1 24 Tf 72 720 Td (" + text.encode("latin-1") + b") Tj ET"
    objs.append(
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream"
    )
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n" + f"0 {len(objs) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += b"trailer\n<< /Size " + str(len(objs) + 1).encode() + b" /Root 1 0 R >>\n"
    out += b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    return bytes(out)


def _make_docx(text: str) -> bytes:
    from docx import Document

    buf = io.BytesIO()
    doc = Document()
    doc.add_paragraph(text)
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(*, long_text: str | None = None, image_only: bool = False) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    if image_only:
        # Minimal valid PNG. Image bytes must never be treated as extracted text.
        png = io.BytesIO(
            bytes.fromhex(
                "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
                "0000000d49444154789c63606060f80f0001040100b51c0c020000000049454e44"
                "ae426082"
            )
        )
        first.shapes.add_picture(png, Inches(1), Inches(1))
    else:
        box = first.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text = "第一页标题"
        table = first.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(1)).table
        table.cell(0, 0).text = "指标"
        table.cell(0, 1).text = "结果"
        table.cell(1, 0).text = "收入"
        table.cell(1, 1).text = "增长"
        group = first.shapes.add_group_shape()
        grouped_box = group.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
        grouped_box.text = "分组形状正文"

        second = presentation.slides.add_slide(presentation.slide_layouts[6])
        second_box = second.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        second_box.text = long_text or "第二页结论"

    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def test_extract_txt():
    r = extract_text("第一行标题\n正文内容".encode(), file_name="a.txt", mime="text/plain")
    assert r.status == "extracted"
    assert "第一行标题" in r.text
    assert r.char_count > 0
    assert r.error_type is None


def test_extract_md():
    r = extract_text(
        b"# Heading\n\n- item\n\n| A | B |\n|---|---|\n| 1 | 2 |\n\n```js\nalert('x')\n```",
        file_name="a.md",
        mime=None,
    )
    assert r.status == "extracted"
    assert "Heading" in r.text
    assert "- item" in r.text
    assert "| A | B |" in r.text
    assert "alert('x')" in r.text


def test_extract_markdown_extension():
    r = extract_text("# 标题\n正文".encode(), file_name="a.markdown", mime="text/markdown")
    assert r.status == "extracted"
    assert "标题" in r.text


def test_extract_md_invalid_encoding_no_500():
    r = extract_text(b"# ok\n\xff\xfe<script>x</script>", file_name="bad.md", mime=None)
    assert r.status == "extracted"
    assert "<script>x</script>" in r.text


def test_extract_txt_non_utf8_robust():
    # 非 UTF-8 字节稳健回退（errors=replace），不应崩溃 / 不应判为 failed。
    r = extract_text(b"\xff\xfe abc", file_name="a.txt", mime="text/plain")
    assert r.status == "extracted"


def test_extract_pdf():
    r = extract_text(_make_pdf("Hello PDF Extract"), file_name="a.pdf", mime="application/pdf")
    assert r.status == "extracted"
    assert "Hello PDF Extract" in r.text


def test_extract_docx():
    r = extract_text(
        _make_docx("供应链优化交付报告"),
        file_name="a.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert r.status == "extracted"
    assert "供应链优化交付报告" in r.text


def test_extract_pptx_text_table_group_in_slide_order():
    r = extract_text(
        _make_pptx(),
        file_name="deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    assert r.status == "extracted"
    assert r.error_type is None
    assert "[幻灯片 1]" in r.text
    assert "指标 | 结果" in r.text
    assert "收入 | 增长" in r.text
    assert "分组形状正文" in r.text
    assert "[幻灯片 2]" in r.text
    assert r.text.index("第一页标题") < r.text.index("分组形状正文")
    assert r.text.index("分组形状正文") < r.text.index("第二页结论")


def test_extract_pptx_by_standard_mime_and_respects_limit():
    marker = "PPTX-LONG-CONTENT-"
    r = extract_text(
        _make_pptx(long_text=marker * (MAX_EXTRACT_CHARS // len(marker) + 100)),
        file_name="upload.bin",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    assert r.status == "extracted"
    assert r.text.startswith("[幻灯片 1]")
    assert len(r.text) == MAX_EXTRACT_CHARS
    assert r.char_count == MAX_EXTRACT_CHARS


def test_extract_image_only_pptx_is_empty():
    r = extract_text(_make_pptx(image_only=True), file_name="image-only.pptx", mime=None)
    assert r.status == "empty"
    assert r.error_type == "extraction_empty"
    assert r.text == ""


def test_extract_corrupt_pptx_fails_without_logging_content_or_path(caplog):
    secret = "PPTX-SECRET-CONTENT C:\\Users\\private\\deck.pptx"
    with caplog.at_level(logging.WARNING):
        r = extract_text(secret.encode(), file_name="broken.pptx", mime=None)
    assert r.status == "failed"
    assert r.error_type == "extraction_failed"
    assert secret not in caplog.text
    assert "C:\\Users\\private" not in (r.error_message or "")


def test_extract_legacy_ppt_is_explicitly_unsupported():
    r = extract_text(b"legacy-binary", file_name="legacy.ppt", mime="application/vnd.ms-powerpoint")
    assert r.status == "unsupported"
    assert r.error_type == "extraction_unsupported"
    assert r.error_message == "当前 .ppt 格式暂不支持自动提取（文件已落盘，请人工补全内容）"


def test_extract_unsupported():
    r = extract_text(b"PK\x03\x04binary", file_name="a.xlsx", mime="application/octet-stream")
    assert r.status == "unsupported"
    assert r.error_type == "extraction_unsupported"
    assert r.text == ""


def test_extract_empty_pdf():
    # 含一页但无文本（如扫描件）→ empty。
    r = extract_text(_make_pdf(""), file_name="scan.pdf", mime="application/pdf")
    assert r.status == "empty"
    assert r.error_type == "extraction_empty"


def test_extract_corrupt_pdf_failed():
    r = extract_text(b"not a real pdf at all", file_name="broken.pdf", mime="application/pdf")
    assert r.status == "failed"
    assert r.error_type == "extraction_failed"
    # 错误信息不得包含敏感路径 / 引用。
    assert "internal://" not in (r.error_message or "")
